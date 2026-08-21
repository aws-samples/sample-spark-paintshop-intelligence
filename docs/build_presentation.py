"""Generate SPARK Paint-shop Intelligence PowerPoint — AWS branded."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── AWS Colour Palette ───────────────────────────────────────────────────────
AWS_SQUID   = RGBColor(0x23, 0x2F, 0x3E)  # AWS Squid Ink — primary background
AWS_ORANGE  = RGBColor(0xFF, 0x99, 0x00)  # AWS Orange — primary accent
AWS_EMBER   = RGBColor(0xEC, 0x72, 0x11)  # AWS Ember — secondary orange
AWS_BLUE    = RGBColor(0x00, 0x73, 0xBB)  # AWS Blue
AWS_GREEN   = RGBColor(0x1A, 0x9C, 0x3E)  # AWS Green
AWS_PURPLE  = RGBColor(0x71, 0x57, 0xFF)  # AWS AI/ML purple
AWS_TEAL    = RGBColor(0x01, 0xA8, 0x8D)  # AWS Teal
MID_DARK    = RGBColor(0x37, 0x47, 0x5A)  # Card / header background
DARK_CARD   = RGBColor(0x1B, 0x26, 0x33)  # Darker card
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xD5, 0xDB, 0xDB)
DIM_GREY    = RGBColor(0x7F, 0x8C, 0x8D)
RED         = RGBColor(0xEF, 0x47, 0x6F)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color=AWS_SQUID):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        font_size=Pt(14), color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = font_size
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    return txb


def rounded_box(slide, left, top, width, height, fill_color, radius=0.1):
    sp = slide.shapes.add_shape(
        5,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    sp.line.fill.background()
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def aws_header(slide, title, subtitle=None):
    """AWS-styled slide header: orange top stripe + squid-ink band + title."""
    box(slide, 0, 0, 13.33, 0.07, fill_color=AWS_ORANGE)
    box(slide, 0, 0.07, 13.33, 1.03, fill_color=MID_DARK)
    txt(slide, title, 0.35, 0.12, 9.5, 0.6,
        font_size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.72, 10.0, 0.32,
            font_size=Pt(11), color=LIGHT_GREY, italic=True)
    txt(slide, "Amazon Web Services", 10.3, 0.18, 2.7, 0.35,
        font_size=Pt(10), color=AWS_ORANGE, bold=True, align=PP_ALIGN.RIGHT)


def aws_footer(slide, note=""):
    box(slide, 0, 7.25, 13.33, 0.17, fill_color=MID_DARK)
    footer_txt = "NAMER AutoMFG BM2026  \u2022  Team 2  \u2022  SPARK"
    if note:
        footer_txt += "  \u2022  " + note
    txt(slide, footer_txt, 0.3, 7.26, 12.73, 0.15,
        font_size=Pt(8), color=DIM_GREY, align=PP_ALIGN.CENTER)
    box(slide, 0, 7.42, 13.33, 0.08, fill_color=AWS_ORANGE)


def divider(slide, y):
    line = slide.shapes.add_shape(
        1, Inches(0.3), Inches(y), Inches(12.73), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = AWS_ORANGE
    line.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide, AWS_SQUID)

# Top & bottom AWS orange stripes
box(slide, 0, 0, 13.33, 0.09, fill_color=AWS_ORANGE)
box(slide, 0, 7.41, 13.33, 0.09, fill_color=AWS_ORANGE)

# AWS label top-right
txt(slide, "Amazon Web Services", 10.0, 0.15, 3.0, 0.38,
    font_size=Pt(11), color=AWS_ORANGE, bold=True, align=PP_ALIGN.RIGHT)

# Central card
rounded_box(slide, 1.4, 1.1, 10.53, 5.3, MID_DARK)

# SPARK title
txt(slide, "SPARK", 2.0, 1.45, 9.33, 1.1,
    font_size=Pt(64), bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)
txt(slide, "Smart Paint-shop Anomaly Response & Knowledge",
    2.0, 2.65, 9.33, 0.65,
    font_size=Pt(20), color=WHITE, align=PP_ALIGN.CENTER, bold=True)

divider(slide, 3.5)

txt(slide, "Real-time anomaly detection  \u2022  Autonomous rescheduling  \u2022  Root-cause analysis",
    2.0, 3.65, 9.33, 0.45,
    font_size=Pt(13), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Service tags
tags = [
    ("Kinesis", AWS_ORANGE),
    ("SageMaker MCE", AWS_PURPLE),
    ("Bedrock AgentCore", AWS_EMBER),
    ("Neptune", AWS_BLUE),
    ("Step Functions", AWS_TEAL),
    ("CDK", AWS_GREEN),
]
x = 1.85
for tag, col in tags:
    w = len(tag) * 0.115 + 0.45
    rounded_box(slide, x, 4.55, w, 0.40, col)
    txt(slide, tag, x + 0.06, 4.57, w - 0.12, 0.36,
        font_size=Pt(10), color=AWS_SQUID, bold=True, align=PP_ALIGN.CENTER)
    x += w + 0.14

# Footer
txt(slide, "NAMER AutoMFG BM2026  \u2022  Team 2  \u2022  2026",
    0, 6.85, 13.33, 0.38,
    font_size=Pt(10), color=DIM_GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Architecture
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "System Architecture",
           "End-to-end AWS cloud-native pipeline — all components deployed")

cols = [
    ("DATA LAYER",      AWS_ORANGE,  0.25,
     ["IoT Simulator\n(Lambda, 1-min cron)",
      "Kinesis Data Stream\n(12 shards, 24h)",
      "Kinesis Firehose\n(Parquet \u2192 S3)"]),
    ("ML PROCESSING",   AWS_PURPLE,  3.6,
     ["Stream Processor\n(Lambda consumer)",
      "SageMaker MCE\n(IF + LSTM + XGBoost\nInService \u2713)",
      "EventBridge\nTankAnomalyDetected"]),
    ("ORCHESTRATION",   AWS_EMBER,   6.95,
     ["Step Functions\n(6-step workflow)",
      "MPS AgentCore\n(reschedule jobs)",
      "RCA AgentCore\n(root-cause report)"]),
    ("STORAGE",         AWS_GREEN,   10.3,
     ["DynamoDB\ntank-status, jobs,\nrca-reports",
      "Neptune KG\n(fault graph,\nGremlin)",
      "S3 + Bedrock KB\n(SOP documents)"]),
]

for col_title, col_color, col_x, items in cols:
    rounded_box(slide, col_x, 1.22, 3.0, 0.42, col_color)
    txt(slide, col_title, col_x, 1.24, 3.0, 0.38,
        font_size=Pt(10), bold=True, color=AWS_SQUID, align=PP_ALIGN.CENTER)
    for i, item in enumerate(items):
        iy = 1.82 + i * 1.6
        rounded_box(slide, col_x + 0.08, iy, 2.84, 1.38, DARK_CARD)
        box(slide, col_x + 0.08, iy, 0.06, 1.38, fill_color=col_color)
        txt(slide, item, col_x + 0.22, iy + 0.18, 2.6, 1.05,
            font_size=Pt(9.5), color=WHITE)

for ax in [3.22, 6.57, 9.92]:
    txt(slide, "\u25ba", ax, 3.3, 0.4, 0.5,
        font_size=Pt(18), color=AWS_ORANGE, align=PP_ALIGN.CENTER)

aws_footer(slide, "CloudFront \u2022 API GW (REST + WebSocket) \u2022 Cognito \u2022 SSM \u2022 IAM \u2022 us-east-1")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Flow
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "Data Flow",
           "From sensor reading to rescheduling decision — fully automated")

steps = [
    (AWS_ORANGE,  "1", "Tank Sensor\nReading",        "12 tanks\n1-min interval"),
    (AWS_TEAL,    "2", "Kinesis\nStream",              "12 shards\n24h retention"),
    (AWS_PURPLE,  "3", "SageMaker\nMCE",               "IF + LSTM +\nXGBoost"),
    (AWS_EMBER,   "4", "EventBridge",                  "TankAnomalyDetected\n(threshold 0.7)"),
    (AWS_BLUE,    "5", "Step Functions\nWorkflow",     "6-step\norchestration"),
]

box_w = 1.85
box_h = 1.65
start_x = 0.28
y = 2.35

for i, (color, num, title, sub) in enumerate(steps):
    x = start_x + i * (box_w + 0.42)
    rounded_box(slide, x, y, box_w, box_h, MID_DARK)
    rounded_box(slide, x + 0.05, y + 0.05, 0.46, 0.46, color)
    txt(slide, num, x + 0.05, y + 0.06, 0.46, 0.40,
        font_size=Pt(14), bold=True, color=AWS_SQUID, align=PP_ALIGN.CENTER)
    txt(slide, title, x + 0.05, y + 0.58, box_w - 0.1, 0.62,
        font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sub, x + 0.05, y + 1.2, box_w - 0.1, 0.44,
        font_size=Pt(9), color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(slide, "\u25ba", x + box_w + 0.07, y + 0.65, 0.35, 0.4,
            font_size=Pt(14), color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Two parallel agent branches
branch_y = 4.35
box(slide, 6.72, y + box_h, 0.02, branch_y - (y + box_h), fill_color=AWS_ORANGE)

for bx, color, agent, tools, output in [
    (0.4,  AWS_EMBER,  "MPS Agent (Bedrock AgentCore)",
     "get_affected_jobs \u2192 get_line_status\n\u2192 compute_reschedule \u2192 apply_schedule",
     "projected_jph \u2022 fbo_delay_mins \u2022 job assignments \u2192 DynamoDB"),
    (6.9,  AWS_GREEN,  "RCA Agent (Bedrock AgentCore)",
     "get_sensor_history \u2192 get_fault_context (Neptune)\n\u2192 get_maintenance_record \u2192 write_rca_report",
     "root_cause \u2022 severity \u2022 recommendation \u2192 rca-reports table"),
]:
    rounded_box(slide, bx, branch_y, 5.9, 2.65, DARK_CARD)
    box(slide, bx, branch_y, 5.9, 0.38, fill_color=color)
    txt(slide, agent, bx + 0.12, branch_y + 0.04, 5.66, 0.32,
        font_size=Pt(11), bold=True, color=AWS_SQUID)
    txt(slide, "Tool chain:", bx + 0.12, branch_y + 0.5, 1.1, 0.25,
        font_size=Pt(9), color=LIGHT_GREY, bold=True)
    txt(slide, tools, bx + 0.12, branch_y + 0.78, 5.66, 0.65,
        font_size=Pt(9), color=WHITE)
    txt(slide, "Output:", bx + 0.12, branch_y + 1.5, 0.85, 0.25,
        font_size=Pt(9), color=LIGHT_GREY, bold=True)
    txt(slide, output, bx + 0.12, branch_y + 1.78, 5.66, 0.55,
        font_size=Pt(9), color=AWS_GREEN if color == AWS_GREEN else AWS_ORANGE)

# Branch arrows
txt(slide, "\u25bc", 4.0, branch_y - 0.42, 0.4, 0.38,
    font_size=Pt(14), color=AWS_EMBER, align=PP_ALIGN.CENTER)
txt(slide, "\u25bc", 9.55, branch_y - 0.42, 0.4, 0.38,
    font_size=Pt(14), color=AWS_GREEN, align=PP_ALIGN.CENTER)

aws_footer(slide)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tank Layout & Process Line
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "Tank Layout — LINE-1",
           "8 pre-treatment + 4 e-coat tanks \u2022 Car body travels left \u2192 right")

sections = [
    ("PRE-CLEAN",  AWS_ORANGE,  0.25,  2.50, ["PT-01\nAlkaline\nDegreaser", "PT-02\nAlkaline\nDegreaser"]),
    ("RINSE",      AWS_TEAL,    2.85,  5.40, ["PT-03\nRinse 1",             "PT-04\nRinse 2"]),
    ("PHOSPHATE",  AWS_GREEN,   5.55,  9.30, ["PT-05\nTi Pre-\nactivation", "PT-06\nZinc\nPhosphate", "PT-07\nRinse 3"]),
    ("SEAL/ED",    AWS_EMBER,   9.45, 11.60, ["PT-08\nPassivation",         "ED-01\nE-Coat Bath"]),
    ("ED RINSE",   AWS_BLUE,   11.75, 13.10, ["ED-02\nED Rinse 1",          "ED-03\nED Rinse 2",  "ED-04\nDI Rinse"]),
]

DARK_TINTS = [
    RGBColor(0x1A, 0x18, 0x00),
    RGBColor(0x00, 0x1F, 0x1C),
    RGBColor(0x01, 0x1E, 0x10),
    RGBColor(0x24, 0x14, 0x00),
    RGBColor(0x00, 0x16, 0x2A),
]

for idx, (sec_name, sec_color, bx1, bx2, tanks) in enumerate(sections):
    box(slide, bx1, 1.2, bx2 - bx1, 5.55, fill_color=DARK_TINTS[idx])
    rounded_box(slide, bx1 + 0.05, 1.25, bx2 - bx1 - 0.1, 0.35, sec_color)
    txt(slide, sec_name, bx1 + 0.05, 1.27, bx2 - bx1 - 0.1, 0.3,
        font_size=Pt(9), bold=True, color=AWS_SQUID, align=PP_ALIGN.CENTER)
    n  = len(tanks)
    tw = min((bx2 - bx1 - 0.2) / n - 0.05, 1.2)
    for j, label in enumerate(tanks):
        tx = bx1 + 0.1 + j * ((bx2 - bx1 - 0.15) / n)
        rounded_box(slide, tx, 1.75, tw, 2.1, MID_DARK)
        box(slide, tx, 1.75, tw, 0.07, fill_color=sec_color)
        txt(slide, label, tx, 1.83, tw, 1.95,
            font_size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER, bold=True)

# Flow arrow
box(slide, 0.25, 4.1, 12.83, 0.1, fill_color=AWS_ORANGE)
txt(slide, "\u25ba", 12.95, 3.95, 0.25, 0.38,
    font_size=Pt(16), color=AWS_ORANGE)
txt(slide, "BODY-IN-WHITE INPUT", 0.3, 4.28, 3.8, 0.28,
    font_size=Pt(9), color=AWS_ORANGE, bold=True)
txt(slide, "COATED BODY OUTPUT \u2192 OVEN", 9.2, 4.28, 4.0, 0.28,
    font_size=Pt(9), color=AWS_GREEN, bold=True)

# Fault monitoring legend
txt(slide, "All 12 tanks monitored by SageMaker MCE  \u2022  9 fault types classified  \u2022  "
          "Alerts via EventBridge",
    0.25, 6.7, 12.83, 0.3,
    font_size=Pt(9), color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

aws_footer(slide)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ML Models
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "ML Models",
           "Three-model ensemble — SageMaker Multi-Container Endpoint (InService)")

models = [
    (AWS_ORANGE, "Isolation\nForest",        "if_score",
     "Unsupervised anomaly detection\nOutputs contamination score 0-1\nSklearn 1.2 \u2022 ml.m5.xlarge",
     "Baseline anomaly flagging\nper sensor reading"),
    (AWS_PURPLE, "LSTM\nAutoencoder",        "lstm_score",
     "Temporal sequence learning\nReconstruction error as signal\nPyTorch 2.0 \u2022 ml.m5.2xlarge",
     "Drift detection over time\nCatches gradual degradation"),
    (AWS_GREEN,  "XGBoost\nFault Classifier","fault_type",
     "9-class fault classification\nReturns label + confidence\nXGBoost 1.7 \u2022 ml.m5.xlarge",
     "Names the fault type:\nzinc_depletion, ph_drift, etc."),
]

for i, (color, name, field, desc, value) in enumerate(models):
    cx = 0.45 + i * 4.2
    rounded_box(slide, cx, 1.3, 3.85, 4.15, MID_DARK)
    rounded_box(slide, cx, 1.3, 3.85, 0.62, color)
    txt(slide, name, cx, 1.35, 3.85, 0.58,
        font_size=Pt(18), bold=True, color=AWS_SQUID, align=PP_ALIGN.CENTER)
    txt(slide, "Output: " + field, cx + 0.12, 2.02, 3.6, 0.32,
        font_size=Pt(10), color=color, bold=True)
    txt(slide, desc, cx + 0.12, 2.38, 3.6, 1.1,
        font_size=Pt(10), color=WHITE)
    txt(slide, "Value:", cx + 0.12, 3.52, 0.75, 0.25,
        font_size=Pt(9), color=LIGHT_GREY, bold=True)
    txt(slide, value, cx + 0.12, 3.78, 3.6, 0.55,
        font_size=Pt(9), color=LIGHT_GREY, italic=True)
    # Status badge
    rounded_box(slide, cx + 0.25, 4.75, 3.35, 0.38,
                RGBColor(0x05, 0x3A, 0x1E))
    txt(slide, "\u2713  DEPLOYED — InService", cx + 0.25, 4.77, 3.35, 0.34,
        font_size=Pt(10), color=AWS_GREEN, bold=True, align=PP_ALIGN.CENTER)

# Endpoint summary bar
rounded_box(slide, 0.45, 5.35, 12.43, 0.85, DARK_CARD)
txt(slide, "Endpoint: paintshop-anomaly-endpoint", 0.65, 5.42, 4.5, 0.28,
    font_size=Pt(11), bold=True, color=AWS_ORANGE)
facts = [
    "InferenceExecutionConfig: Direct (per-container targeting)",
    "Instance: ml.m5.xlarge  \u2022  All 3 containers on same host",
    "Training data: 6 months \u00d7 12 tanks \u2022 ~960k labelled rows",
    "Pipeline: PaintShopAnomalyPipeline  \u2022  SageMaker Model Registry",
]
for j, f in enumerate(facts):
    col = j // 2
    row = j % 2
    txt(slide, "\u2022 " + f, 0.65 + col * 6.1, 5.72 + row * 0.27, 5.9, 0.26,
        font_size=Pt(9), color=LIGHT_GREY)

aws_footer(slide)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Fault Types by Tank
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "Fault Types by Tank",
           "9 fault signatures \u2022 All classified by XGBoost \u2022 All seeded in Neptune knowledge graph")

defects = [
    ("acid_drift",            RED,         "PT-06",              "free_acid_pts HIGH",        "Coating too thick / brittle"),
    ("zinc_depletion",        AWS_ORANGE,  "PT-06",              "zinc_g_per_l LOW",           "Bare spots — no coverage"),
    ("accelerator_depletion", AWS_TEAL,    "PT-06",              "accelerator_pts LOW",        "Slow / uneven crystal growth"),
    ("titanium_depletion",    AWS_BLUE,    "PT-05",              "titanium_ppm LOW",           "Coarse crystals, poor adhesion"),
    ("alkalinity_depletion",  AWS_GREEN,   "PT-01 / PT-02",      "free_alkalinity LOW",        "Oil carry-over, poor degreasing"),
    ("rinse_contamination",   AWS_PURPLE,  "PT-03/04/07\nED-02/03/04", "conductivity HIGH",   "Chemical drag-over, adhesion loss"),
    ("temperature_creep",     AWS_EMBER,   "ED-01",              "temperature_c HIGH",         "Paint film defects, blistering"),
    ("meq_acid_buildup",      AWS_ORANGE,  "ED-01",              "meq_acid HIGH",              "Bath acidification, film reject"),
    ("ph_drift",              AWS_TEAL,    "PT-08 / ED-01",      "ph HIGH",                    "Coating weight variance"),
]

cols_per_row = 3
card_w = 4.12
card_h = 1.22

for idx, (fault, color, tank, sensor, impact) in enumerate(defects):
    row = idx // cols_per_row
    col = idx % cols_per_row
    cx = 0.28 + col * (card_w + 0.16)
    cy = 1.28 + row * (card_h + 0.13)
    rounded_box(slide, cx, cy, card_w, card_h, MID_DARK)
    box(slide, cx, cy, 0.07, card_h, fill_color=color)
    txt(slide, fault.replace("_", " ").upper(), cx + 0.14, cy + 0.05, 2.5, 0.3,
        font_size=Pt(9), bold=True, color=color)
    rounded_box(slide, cx + 2.72, cy + 0.04, 1.3, 0.28, color)
    txt(slide, tank, cx + 2.72, cy + 0.06, 1.3, 0.25,
        font_size=Pt(7.5), bold=True, color=AWS_SQUID, align=PP_ALIGN.CENTER)
    txt(slide, "Sensor: " + sensor, cx + 0.14, cy + 0.40, 3.88, 0.28,
        font_size=Pt(9), color=LIGHT_GREY)
    txt(slide, "Impact: " + impact, cx + 0.14, cy + 0.70, 3.88, 0.44,
        font_size=Pt(8.5), color=WHITE, italic=True)

aws_footer(slide, "Neptune KG: 33 FaultType vertices \u2022 33 SOPs \u2022 causal chain edges")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI Agents
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "AI Agents — Bedrock AgentCore",
           "Two autonomous agents \u2022 Strands framework \u2022 Claude Haiku 4.5 \u2022 ~55s end-to-end")

for ax, color, name, role, tools, output_fields in [
    (0.3, AWS_EMBER, "MPS Agent",
     "Master Production Schedule Supervisor\nReschedules affected jobs when a tank degrades",
     [
         ("get_affected_jobs",  "Fetches IN_PROGRESS + QUEUED jobs for the offline tank"),
         ("get_line_status",    "Identifies healthy alternative tanks on the same line"),
         ("compute_reschedule", "Runs constraints solver — returns optimal assignment plan"),
         ("apply_schedule",     "Commits new assignments to production-jobs DynamoDB table"),
     ],
     [("projected_jph", AWS_ORANGE), ("fbo_delay_mins", AWS_TEAL),
      ("rerouted_count", AWS_GREEN), ("held_count", LIGHT_GREY)]),
    (6.85, AWS_GREEN, "RCA Agent",
     "Root-Cause Analysis Investigator\nIdentifies why a fault occurred and recommends remediation",
     [
         ("get_sensor_history",    "Retrieves 6-hour telemetry window from DynamoDB"),
         ("get_fault_context",     "Queries Neptune KG — severity, SOP, causal chain"),
         ("get_maintenance_record","Checks maintenance history for the tank"),
         ("write_rca_report",      "Persists full RCA report to rca-reports DynamoDB table"),
     ],
     [("severity", RED), ("root_cause", AWS_ORANGE), ("recommendation", LIGHT_GREY),
      ("report_id", DIM_GREY)]),
]:
    card_w = 6.2
    rounded_box(slide, ax, 1.22, card_w, 5.85, MID_DARK)
    box(slide, ax, 1.22, card_w, 0.45, fill_color=color)
    txt(slide, name, ax + 0.15, 1.25, 4.0, 0.40,
        font_size=Pt(15), bold=True, color=AWS_SQUID)
    # Model badge
    rounded_box(slide, ax + 4.3, 1.25, 1.75, 0.35, AWS_SQUID)
    txt(slide, "claude-3.5-haiku", ax + 4.35, 1.27, 1.65, 0.30,
        font_size=Pt(8), color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, role, ax + 0.15, 1.75, card_w - 0.3, 0.52,
        font_size=Pt(9.5), color=LIGHT_GREY, italic=True)

    txt(slide, "TOOL CHAIN", ax + 0.15, 2.35, 2.0, 0.25,
        font_size=Pt(8), color=color, bold=True)
    for ti, (tname, tdesc) in enumerate(tools):
        ty = 2.62 + ti * 0.65
        rounded_box(slide, ax + 0.15, ty, card_w - 0.3, 0.56, DARK_CARD)
        box(slide, ax + 0.15, ty, 0.05, 0.56, fill_color=color)
        txt(slide, tname, ax + 0.26, ty + 0.04, 2.1, 0.25,
            font_size=Pt(9), bold=True, color=color)
        txt(slide, tdesc, ax + 0.26, ty + 0.28, card_w - 0.55, 0.24,
            font_size=Pt(8.5), color=LIGHT_GREY)

    txt(slide, "OUTPUT FIELDS", ax + 0.15, 5.27, 2.5, 0.25,
        font_size=Pt(8), color=color, bold=True)
    fx = ax + 0.15
    for fname, fcolor in output_fields:
        fw = len(fname) * 0.1 + 0.35
        rounded_box(slide, fx, 5.52, fw, 0.3, DARK_CARD)
        txt(slide, fname, fx + 0.06, 5.54, fw - 0.1, 0.25,
            font_size=Pt(8), color=fcolor, bold=True)
        fx += fw + 0.1

# Arrow between agents
txt(slide, "\u2194", 6.3, 3.5, 0.45, 0.5,
    font_size=Pt(20), color=DIM_GREY, align=PP_ALIGN.CENTER)
txt(slide, "parallel\nexecution", 6.12, 4.05, 0.85, 0.45,
    font_size=Pt(7.5), color=DIM_GREY, align=PP_ALIGN.CENTER, italic=True)

aws_footer(slide, "AgentCore Gateway \u2022 Cognito M2M OAuth \u2022 Lambda tool targets")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Deployed AWS Infrastructure
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "Deployed AWS Infrastructure",
           "9 CDK stacks \u2022 All CREATE_COMPLETE \u2022 us-east-1")

stacks = [
    ("PaintShopStorage",    AWS_ORANGE,
     ["S3: deployment-configured ML bucket", "Versioned + lifecycle rules", "Audit log bucket"]),
    ("PaintShopIam",        AWS_BLUE,
     ["8 IAM roles", "Least-privilege policies", "Cross-stack references"]),
    ("PaintShopSageMaker",  AWS_PURPLE,
     ["MCE Endpoint (InService)", "3 training containers", "Model Package Group"]),
    ("PaintShopIngestion",  AWS_TEAL,
     ["Kinesis (12 shards)", "Firehose \u2192 Parquet/S3", "Glue Catalog"]),
    ("PaintShopNeptune",    AWS_EMBER,
     ["Neptune Serverless 2-8 NCU", "33 FaultType vertices", "Gremlin endpoint"]),
    ("PaintShopBedrock",    AWS_GREEN,
     ["Bedrock Knowledge Base", "SOP documents (OpenSearch)", "MPS + RCA invoker Lambdas"]),
    ("PaintShopScheduling", AWS_ORANGE,
     ["3 DynamoDB tables", "Step Functions (6 steps)", "EventBridge rule \u2192 SFN"]),
    ("PaintShopApi",        AWS_BLUE,
     ["REST API + WebSocket API", "Cognito User Pool", "Lambda URL (agent stream)"]),
    ("PaintShopFrontend",   AWS_PURPLE,
     ["CloudFront distribution", "S3 static site bucket", "WAF web ACL"]),
]

card_w2 = 4.05
card_h2 = 1.75
for i, (name, color, bullets) in enumerate(stacks):
    row = i // 3
    col = i % 3
    cx = 0.22 + col * (card_w2 + 0.27)
    cy = 1.28 + row * (card_h2 + 0.18)
    rounded_box(slide, cx, cy, card_w2, card_h2, MID_DARK)
    box(slide, cx, cy, card_w2, 0.4, fill_color=color)
    txt(slide, "\u2713", cx + 0.08, cy + 0.05, 0.3, 0.32,
        font_size=Pt(13), bold=True, color=AWS_SQUID, align=PP_ALIGN.CENTER)
    txt(slide, name, cx + 0.38, cy + 0.06, card_w2 - 0.46, 0.28,
        font_size=Pt(10), bold=True, color=AWS_SQUID)
    for j, b in enumerate(bullets):
        txt(slide, "\u2022 " + b, cx + 0.1, cy + 0.5 + j * 0.38, card_w2 - 0.2, 0.36,
            font_size=Pt(9), color=LIGHT_GREY)

aws_footer(slide, "CDK v2 (Python) \u2022 Infrastructure as Code \u2022 All stacks in dependency order")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Live Demo
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
aws_header(slide, "Live Demo",
           "Fault injection \u2192 ML detection \u2192 autonomous agent response \u2022 ~90 seconds end-to-end")

steps_demo = [
    (AWS_ORANGE, "1", "Open Dashboard",
     "Navigate to CloudFront URL\nLog in with Cognito credentials"),
    (AWS_TEAL,   "2", "Inject Fault",
     "Click \"Inject [fault]\" on any tank\nFault drifts in over ~90 seconds via SSM"),
    (AWS_PURPLE, "3", "ML Detection",
     "SageMaker MCE scores readings\nIF + LSTM + XGBoost classify fault type"),
    (AWS_EMBER,  "4", "Tank Degrades",
     "Status flips to DEGRADED\nEventBridge fires TankAnomalyDetected"),
    (AWS_GREEN,  "5", "Agents Activate",
     "Step Functions triggers both agents\nMPS reschedules, RCA investigates (~55s)"),
    (AWS_BLUE,   "6", "View Results",
     "Dashboard shows projected JPH\nRCA report with root cause + recommendation"),
]

bw = 1.88
bh = 2.05
sx = 0.3
sy = 1.45

for i, (color, num, title, desc) in enumerate(steps_demo):
    x = sx + i * (bw + 0.26)
    rounded_box(slide, x, sy, bw, bh, MID_DARK)
    rounded_box(slide, x + 0.05, sy + 0.05, 0.48, 0.48, color)
    txt(slide, num, x + 0.05, sy + 0.06, 0.48, 0.40,
        font_size=Pt(14), bold=True, color=AWS_SQUID, align=PP_ALIGN.CENTER)
    txt(slide, title, x + 0.06, sy + 0.60, bw - 0.12, 0.42,
        font_size=Pt(11), bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, desc, x + 0.1, sy + 1.05, bw - 0.2, 0.92,
        font_size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps_demo) - 1:
        txt(slide, "\u25ba", x + bw + 0.03, sy + 0.75, 0.25, 0.38,
            font_size=Pt(13), color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Key metrics bar
rounded_box(slide, 0.3, 3.82, 12.73, 2.95, DARK_CARD)
txt(slide, "Key Metrics", 0.55, 3.9, 3.0, 0.38,
    font_size=Pt(14), bold=True, color=AWS_ORANGE)
divider(slide, 4.35)

metrics = [
    ("12", "Tanks monitored",          AWS_ORANGE),
    ("9",  "Fault types classified",   AWS_PURPLE),
    ("3",  "ML models in ensemble",    AWS_GREEN),
    ("~55s","Agent execution time",    AWS_TEAL),
    ("9",  "CDK stacks deployed",      AWS_BLUE),
    ("33", "Neptune fault vertices",   AWS_EMBER),
]
for j, (val, label, color) in enumerate(metrics):
    mx = 0.55 + j * 2.1
    txt(slide, val, mx, 4.45, 1.8, 0.75,
        font_size=Pt(32), bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label, mx, 5.22, 1.8, 0.35,
        font_size=Pt(9), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

txt(slide, "Reset to Normal after demo — tank returns to online status within 1 minute",
    0.3, 6.55, 12.73, 0.3,
    font_size=Pt(9.5), color=DIM_GREY, align=PP_ALIGN.CENTER, italic=True)

aws_footer(slide)


# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "paint-shop-presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
